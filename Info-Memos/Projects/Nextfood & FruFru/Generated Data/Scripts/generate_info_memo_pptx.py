#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Generate Nextfood & FruFru Info Memo PPTX (4 slides) — V2

Sell-side info memo for PE buyers.
All financial data read from CSVs, CAGR calculated by Python.
KPMG-inspired layout: widescreen 10.83" x 7.50"

V2 changes:
- Complete P&L tables with ALL line items (historic, consolidated, budgets)
- Fixed team: Simiuc = FruFru, Lifebox = Balaceanu, Scarlat. OTOTO not in package (Simiuc's separate business)
- Added market context from research + Poland as CEE leading indicator
- Cleaned up key risks (factual only, no speculative assumptions)
- 4 slides: Brand & Team | Historic Financials | Budget | Highlights & Risks

Sources:
  - Financial CSVs: ../Financial CSVs/
  - Logos: ../../Logos/
  - Research notes: ../../Notes/Research/
"""

import csv
import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ==================== PATHS ====================

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(os.path.dirname(SCRIPT_DIR))
CSV_DIR = os.path.join(os.path.dirname(SCRIPT_DIR), "Financial CSVs")
LOGO_DIR = os.path.join(PROJECT_DIR, "Logos")
OUTPUT_DIR = os.path.dirname(SCRIPT_DIR)

FRUFRU_LOGO = os.path.join(LOGO_DIR, "frufru-logo.png")
LIFEBOX_LOGO = os.path.join(LOGO_DIR, "llifebox-logo.png")
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, "Nextfood_FruFru_Info_Memo.pptx")
V7_LOGO = os.path.join(LOGO_DIR, "v7capital-logo.png")

# ==================== COLORS (V7 Capital Brand Guidelines 2026) ====================

PRIMARY_BLUE = RGBColor(0x0D, 0x1C, 0x43)    # Dull Violet Black — headers, titles
ACCENT_BLUE = RGBColor(0xBC, 0x89, 0x2B)     # Khaki — accent, section headers, table headers
VERY_LIGHT_BLUE = RGBColor(0xF0, 0xF3, 0xE7) # Pale King's Blue A — background highlights
DARK_TEXT = RGBColor(0x0D, 0x1C, 0x43)        # Dull Violet Black — body text
MEDIUM_TEXT = RGBColor(0x0D, 0x1C, 0x43)      # Dull Violet Black — subtitles
LIGHT_TEXT = RGBColor(0x66, 0x66, 0x66)       # footnotes, captions
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_ACCENT = RGBColor(0xBC, 0x89, 0x2B)     # Khaki — positive highlights
ROW_ALT = RGBColor(0xF0, 0xF3, 0xE7)          # Pale King's Blue A — alternating rows

# ==================== EXCHANGE RATES ====================

# RON/EUR average annual exchange rates (from BNR / Romania Ministry of Finance)
EXCHANGE_RATES_FILE = os.path.join(CSV_DIR, "exchange_rates_ron_eur.json")


def load_exchange_rates():
    """Load RON/EUR exchange rates from JSON file"""
    with open(EXCHANGE_RATES_FILE, 'r', encoding='utf-8') as f:
        data = json.load(f)
    # Return only numeric year keys as {year_str: rate}
    return {k: v for k, v in data.items() if not k.startswith('_')}


def to_eur(val_ron, year):
    """Convert RON value to EUR using the year's average exchange rate"""
    rates = EXCHANGE_RATES
    rate = rates.get(str(year))
    if rate is None or rate == 0 or val_ron is None:
        return None
    return val_ron / rate


def fmt_eur(val, millions=True):
    """Format EUR value for display"""
    if val is None:
        return "N/A"
    if millions:
        return f"€{val / 1_000_000:,.1f}M"
    return f"€{val:,.0f}"


def eur(val, year):
    """Convenience: convert RON value to EUR and format for display (millions)"""
    return fmt_eur(to_eur(val, year))


# Load exchange rates at module level (used throughout)
EXCHANGE_RATES = load_exchange_rates()


# ==================== CSV READING ====================

def read_csv(filename):
    filepath = os.path.join(CSV_DIR, filename)
    rows = []
    with open(filepath, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(row)
    return rows


def get_csv_dict(filename):
    rows = read_csv(filename)
    result = {}
    for row in rows:
        if row and row[0]:
            result[row[0].strip()] = row[1:]
    return result


def parse_float(val):
    if not val or val.strip() == '':
        return 0.0
    try:
        return float(val.strip())
    except ValueError:
        return 0.0


def calc_cagr(start_val, end_val, years):
    if start_val <= 0 or end_val <= 0 or years <= 0:
        return None
    return (end_val / start_val) ** (1.0 / years) - 1.0


def fmt_ron(val, millions=True):
    if val is None:
        return "N/A"
    if millions:
        return f"{val / 1_000_000:,.1f}M"
    return f"{val:,.0f}"


def fmt_pct(val):
    if val is None:
        return "N/A"
    return f"{val * 100:+.1f}%"


def fmt_pct_abs(val):
    if val is None:
        return "N/A"
    return f"{val * 100:.1f}%"


def fmt_int(val):
    if val is None or val == 0:
        return "—"
    return f"{val:,.0f}"


# ==================== LOAD ALL FINANCIAL DATA ====================

def load_financials():
    """Load all CSV financial data — ALL line items (identical logic to PDF script)"""
    data = {}

    # --- Historic Lifebox 2020-2023 ---
    hist = get_csv_dict("historic_lifebox_2020_2023.csv")
    data['historic'] = {'years': ['2020', '2021', '2022', '2023']}
    for m in ['Boxes_Produced', 'Net_Revenue', 'Material_Costs', 'Material_Margin',
              'Labor_Direct_Indirect', 'Manufacturing', 'Depreciation_Production',
              'Gross_Margin', 'Transport_Costs', 'Marketing', 'GA',
              'EBITDA_Normalised', 'Net_Profit']:
        data['historic'][m] = [parse_float(v) for v in hist.get(m, ['0'] * 4)]

    # --- Consolidated 2023 vs 2024 ---
    cons = get_csv_dict("consolidat_2023_2024.csv")
    data['consolidated'] = {'years': ['2023', '2024']}
    for m in ['Cantitate_Produsa_KG', 'Boxuri_Produse', 'Net_Revenue', 'Material_Costs',
              'Material_Margin', 'Labor_Direct_Indirect', 'Manufacturing',
              'Depreciation_Production', 'Gross_Margin', 'Transport_Costs',
              'Marketing', 'GA', 'Depreciation_GA', 'EBITDA_Normalised', 'Net_Profit']:
        data['consolidated'][m] = [parse_float(v) for v in cons.get(m, ['0'] * 2)]

    # --- Consolidated 2025 ---
    cons25 = get_csv_dict("consolidat_2025_monthly.csv")
    data['consolidated_2025'] = {}
    for m in ['Sales', 'COGS', 'Food_Cost', 'Packaging', 'Gross_Margin_After_COGS',
              'Labor_Direct_Indirect', 'Manufacturing', 'Depreciation_Production',
              'Delivery', 'Marketing', 'GA', 'EBITDA', 'Net_Profit']:
        vals = cons25.get(m, [])
        data['consolidated_2025'][m] = parse_float(vals[-1]) if vals else 0.0

    # --- Budget 2025 by BU ---
    for bu_name, bu_file in [
        ('frufru', 'bu_frufru_2025_budget.csv'),
        ('lifebox', 'bu_lifebox_2025_budget.csv'),
        ('catering', 'bu_catering_2025_budget.csv'),
    ]:
        bu = get_csv_dict(bu_file)
        bu_data = {}
        for m in ['Sales', 'COGS', 'Food_Cost', 'Packaging', 'Gross_Margin',
                   'Labor_Direct_Indirect', 'Manufacturing', 'Depreciation',
                   'Delivery', 'Marketing', 'GA']:
            vals = bu.get(m, [])
            bu_data[m] = parse_float(vals[-1]) if vals else 0.0
        data[f'bu_{bu_name}_2025'] = bu_data

    # --- Budget 2026 base case ---
    b26 = get_csv_dict("buget_2026_base.csv")
    data['budget_2026'] = {}
    all_b26_metrics = [
        'Sales_Lifebox', 'Sales_FruFru', 'Sales_Catering', 'Total_Sales',
        'Food_Cost', 'Cost_Ambalaj', 'Total_Material_Cost', 'Material_Margin',
        'Salarii_Directe', 'Salarii_Indirecte', 'Total_Salarii',
        'Total_Cost_Productie', 'Gross_Margin', 'Delivery', 'Marketing', 'GA',
        'EBITDA_Normalised'
    ]
    for m in all_b26_metrics:
        vals = b26.get(m, [])
        data['budget_2026'][m + '_H1'] = parse_float(vals[6]) if len(vals) > 6 else 0.0
        data['budget_2026'][m + '_FY'] = parse_float(vals[7]) if len(vals) > 7 else 0.0
    data['budget_2026']['Total_Sales'] = data['budget_2026']['Total_Sales_FY']
    data['budget_2026']['EBITDA_Normalised'] = data['budget_2026']['EBITDA_Normalised_FY']

    # --- Derived metrics ---
    rev_2024 = data['consolidated']['Net_Revenue'][1]
    rev_2025 = data['consolidated_2025']['Sales']
    rev_2026 = data['budget_2026']['Total_Sales']
    data['budget_2025_total'] = rev_2025
    data['cagr_rev_2024_2026'] = calc_cagr(rev_2024, rev_2026, 2)
    rev_2023_c = data['consolidated']['Net_Revenue'][0]
    data['growth_rev_2023_2024'] = (rev_2024 / rev_2023_c - 1.0) if rev_2023_c > 0 else None
    rev_2020 = data['historic']['Net_Revenue'][0]
    rev_2023_h = data['historic']['Net_Revenue'][3]
    data['cagr_historic_2020_2023'] = calc_cagr(rev_2020, rev_2023_h, 3)

    return data


# ==================== PPTX HELPERS ====================

def set_cell_text(cell, text, font_size=8, bold=False, color=None, alignment=PP_ALIGN.CENTER):
    """Set text in a table cell with formatting"""
    if color is None:
        color = DARK_TEXT
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_fill(cell, color):
    """Set background fill for a table cell"""
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=10,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide"""
    if color is None:
        color = DARK_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=9,
                       color=None, bold_prefix=True):
    """Add a text box with bullet points.
    bullets = [("Bold prefix.", "Rest of text"), ...]
    """
    if color is None:
        color = DARK_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (prefix, rest) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        p.space_before = Pt(1)

        bullet_run = p.add_run()
        bullet_run.text = "• "
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.color.rgb = color
        bullet_run.font.name = "Calibri"

        prefix_run = p.add_run()
        prefix_run.text = prefix
        prefix_run.font.size = Pt(font_size)
        prefix_run.font.bold = bold_prefix
        prefix_run.font.color.rgb = color
        prefix_run.font.name = "Calibri"

        if rest:
            rest_run = p.add_run()
            rest_run.text = " " + rest
            rest_run.font.size = Pt(font_size)
            rest_run.font.bold = False
            rest_run.font.color.rgb = color
            rest_run.font.name = "Calibri"

    return txBox


def add_section_banner(slide, left, top, width, text):
    """Add a colored section banner"""
    shape = slide.shapes.add_shape(1, left, top, width, Cm(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(0.3)


def add_table(slide, left, top, width, rows_data, col_widths_cm,
              header_color=None, font_size=7, bold_rows=None, row_height_cm=0.45):
    """Add a formatted table to a slide.
    rows_data = [['Header1', ...], ['val1', ...], ...]
    col_widths_cm = [3, 4, 5, ...] (relative proportions)
    bold_rows = set of 1-based row indices to render bold with highlight
    row_height_cm = height of each row in cm (default 0.45)
    """
    if header_color is None:
        header_color = ACCENT_BLUE
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    row_height = Cm(row_height_cm)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, row_height * n_rows
    )
    table = table_shape.table

    total_cm = sum(col_widths_cm)
    for i, w in enumerate(col_widths_cm):
        table.columns[i].width = int(width * w / total_cm)

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            if row_idx == 0:
                set_cell_text(cell, cell_text, font_size=font_size, bold=True, color=WHITE)
                set_cell_fill(cell, header_color)
            else:
                is_first_col = col_idx == 0
                is_bold = bold_rows and row_idx in bold_rows
                set_cell_text(cell, cell_text, font_size=font_size,
                              bold=is_bold, color=DARK_TEXT,
                              alignment=PP_ALIGN.LEFT if is_first_col else PP_ALIGN.CENTER)
                if is_bold:
                    set_cell_fill(cell, VERY_LIGHT_BLUE)
                elif row_idx % 2 == 0:
                    set_cell_fill(cell, ROW_ALT)
            cell.margin_left = Cm(0.1)
            cell.margin_right = Cm(0.1)
            cell.margin_top = Cm(0.03)
            cell.margin_bottom = Cm(0.03)

    return table_shape


def add_kpi_box(slide, left, top, width, height, value_text, label_text,
                value_color=None):
    """Add a KPI display box"""
    if value_color is None:
        value_color = ACCENT_BLUE
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = VERY_LIGHT_BLUE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(0.5)
    add_textbox(slide, left, top + Cm(0.1), width, Cm(0.7),
                value_text, font_size=14, bold=True, color=value_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Cm(0.75), width, Cm(0.4),
                label_text, font_size=7, bold=False, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)


# ==================== PPTX GENERATION ====================

def create_info_memo_pptx():
    """Create the 4-slide info memo PPTX"""

    fin = load_financials()

    prs = Presentation()
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    margin_left = Cm(1.5)
    margin_right = Cm(1.5)
    content_width = slide_w - margin_left - margin_right

    # ==================== SLIDE 1: BRAND & TEAM ====================

    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # --- Logos ---
    if os.path.exists(FRUFRU_LOGO):
        slide1.shapes.add_picture(FRUFRU_LOGO, margin_left, Cm(0.5), width=Cm(5), height=Cm(1.08))
    if os.path.exists(LIFEBOX_LOGO):
        slide1.shapes.add_picture(LIFEBOX_LOGO, slide_w - margin_right - Cm(3), Cm(0.3),
                                  width=Cm(2.4), height=Cm(2))

    # --- Title ---
    add_textbox(slide1, margin_left, Cm(1.8), content_width, Cm(1.0),
                "Nextfood & FruFru — Information Memorandum",
                font_size=22, bold=True, color=PRIMARY_BLUE)
    add_textbox(slide1, margin_left, Cm(2.7), content_width, Cm(0.4),
                f"Confidential | {datetime.now().strftime('%B %Y')} | Prepared for prospective investors",
                font_size=10, color=MEDIUM_TEXT)

    # --- Line ---
    line = slide1.shapes.add_shape(1, margin_left, Cm(3.2), content_width, Cm(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # --- Brand Overview ---
    add_section_banner(slide1, margin_left, Cm(3.4), content_width, "BRAND OVERVIEW")
    col_w = int(content_width / 2 - Cm(0.2))

    add_textbox(slide1, margin_left, Cm(4.4), col_w, Cm(3.0),
        "FruFru — Founded 2006 by Mihai Simiuc as the first fresh healthy food concept in Romania. "
        "Sold to Unilever (75%) in 2019, bought back January 2024. Relaunched October 2024. "
        "~30+ SKUs: fresh salads, soups, bowls, mains, desserts. "
        "Brand DNA: #farabazaconii — no preservatives, no nonsense. "
        "Distributed via Freshful, Kaufland, Mega Image, Carrefour, lifebox.ro, and OTOTO stores "
        "(Simiuc's separate retail chain — distribution synergies).",
        font_size=8, color=DARK_TEXT)

    add_textbox(slide1, margin_left + col_w + Cm(0.4), Cm(4.4), col_w, Cm(3.0),
        "Lifebox (Nextfood SRL) — Healthy meal subscription service with daily delivery. "
        "Founded by Radu Balaceanu and Florin Scarlat. "
        "7 personalised menu types (OptimBox, Vegan, Sport, Custom, etc.) designed by nutritionists. "
        "Delivery zones: Bucharest, Cluj, Oradea, Budapest (as FrissBox). "
        "Shared production infrastructure with FruFru — the Unilever buyback excluded the factory.",
        font_size=8, color=DARK_TEXT)

    # --- Market context ---
    add_textbox(slide1, margin_left, Cm(7.4), content_width, Cm(1.0),
        "Market Context: Romanian RTE market ~EUR 1.75B (2025), growing 3.7% CAGR. Healthy sub-segment "
        "EUR 250-350M. Romania at 78% of EU average income, converging. "
        "Poland, a comparable CEE market, has seen rapid ready meals adoption — "
        "a trajectory Romania is expected to follow.",
        font_size=7, color=LIGHT_TEXT)

    # --- Entity ---
    entities_data = [
        ['Entity', 'Brands', 'Role'],
        ['Nextfood SRL', 'Lifebox + FruFru', 'Production, meal delivery, packaged retail (#farabazaconii)'],
    ]
    add_table(slide1, margin_left, Cm(8.5), content_width,
              entities_data, [3, 3, 9], font_size=7)

    # --- Ownership structure (from constitutive act) ---
    ownership_data = [
        ['Shareholder', 'Shares', '%'],
        ['Mihai Simiuc', '115,500', '38.5%'],
        ['Vertical Seven Group S.A.', '120,000', '40.0%'],
        ['Scarlat Florin-Ioan', '45,000', '15.0%'],
        ['Radu Balaceanu', '19,500', '6.5%'],
    ]
    add_table(slide1, margin_left, Cm(9.5), content_width,
              ownership_data, [7, 3, 2], font_size=7)

    # --- Team ---
    # Ownership table: 9.5 + 5×0.45 = ~11.75cm
    add_section_banner(slide1, margin_left, Cm(12.0), content_width, "TEAM")

    add_textbox(slide1, margin_left, Cm(13.0), content_width, Cm(1.2),
        "FruFru — Mihai Simiuc (Founder & CEO). Serial entrepreneur, HEC MBA. "
        "Founded FruFru (2006), grew to RON 34.4M revenue, sold to Unilever (2019), bought back 2024. "
        "Also operates OTOTO retail chain (separate business, distribution synergies). "
        "Senior Adviser at Comitis Capital (Frankfurt PE).",
        font_size=8, color=DARK_TEXT)

    add_textbox(slide1, margin_left, Cm(14.2), content_width, Cm(1.0),
        "Lifebox — Radu Balaceanu, Florin Scarlat (Co-founders). "
        "Team: Dr. Anamaria Iulian (Nutritionist), Chef Alex Cirtu (Head Chef). "
        "~60 kitchen staff, 1,500 sqm production facility (APACA, Bucharest). ~1,000 boxes/day.",
        font_size=8, color=DARK_TEXT)

    # --- Timeline ---
    timeline_data = [
        ['Year', 'Event'],
        ['2006', 'FruFru & Urban Monkey founded by Simiuc'],
        ['2017', 'Lifebox founded by Balaceanu, Scarlat'],
        ['2019', 'Unilever acquires 75% of Good People SA'],
        ['2024 Jan', 'Simiuc buys back FruFru & Urban Monkey from Unilever'],
        ['2024 Oct', 'FruFru relaunched via retail (Mega Image, Carrefour, Freshful) and OTOTO stores'],
    ]
    # Team ends at ~15.2cm, timeline 6 rows × 0.45 = 2.7cm → ends ~18.2cm (tight but fits)
    add_table(slide1, margin_left, Cm(15.4), content_width,
              timeline_data, [2, 15], header_color=PRIMARY_BLUE, font_size=7)

    # --- Footer ---
    add_textbox(slide1, margin_left, Cm(17.7), content_width, Cm(0.4),
                "CONFIDENTIAL",
                font_size=6, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)

    # ==================== SLIDE 2: HISTORIC & CURRENT FINANCIALS ====================

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_banner(slide2, margin_left, Cm(0.3), content_width,
                       "FINANCIAL PERFORMANCE — HISTORIC & CURRENT")

    # --- KPI boxes ---
    rev_2024 = fin['consolidated']['Net_Revenue'][1]
    rev_2025_bgt = fin['budget_2025_total']
    rev_2026_bgt = fin['budget_2026']['Total_Sales']
    ebitda_2026 = fin['budget_2026']['EBITDA_Normalised']
    ebitda_margin_2026 = ebitda_2026 / rev_2026_bgt if rev_2026_bgt > 0 else None

    cagr_rev = fin['cagr_rev_2024_2026']

    kpi_w = int(content_width / 5 - Cm(0.2))
    kpi_y = Cm(1.3)
    kpi_gap = Cm(0.2)
    add_kpi_box(slide2, margin_left, kpi_y, kpi_w, Cm(1.1),
                eur(rev_2024, 2024), "Revenue 2024")
    add_kpi_box(slide2, margin_left + 1 * (kpi_w + kpi_gap), kpi_y, kpi_w, Cm(1.1),
                eur(rev_2025_bgt, 2025), "Budget 2025")
    add_kpi_box(slide2, margin_left + 2 * (kpi_w + kpi_gap), kpi_y, kpi_w, Cm(1.1),
                eur(rev_2026_bgt, 2026), "Budget 2026")
    add_kpi_box(slide2, margin_left + 3 * (kpi_w + kpi_gap), kpi_y, kpi_w, Cm(1.1),
                eur(ebitda_2026, 2026), "EBITDA 2026E", value_color=GREEN_ACCENT)
    add_kpi_box(slide2, margin_left + 4 * (kpi_w + kpi_gap), kpi_y, kpi_w, Cm(1.1),
                fmt_pct(cagr_rev), "CAGR 24→26E", value_color=GREEN_ACCENT)

    # --- Single Consolidated P&L: 2023, 2024, 2025B ---
    add_textbox(slide2, margin_left, Cm(2.6), content_width, Cm(0.4),
                "Consolidated P&L — Lifebox + Catering (2023–2024) | + FruFru (2025B)",
                font_size=9, bold=True, color=PRIMARY_BLUE)

    c = fin['consolidated']
    c25 = fin['consolidated_2025']
    growth_rev = fin['growth_rev_2023_2024']

    merged_data = [
        ['EUR', '2023', '2024', '2025 B'],
        ['Net Revenue / Sales',
         eur(c['Net_Revenue'][0], 2023), eur(c['Net_Revenue'][1], 2024), eur(c25['Sales'], 2025)],
        ['  Material Costs / COGS',
         eur(c['Material_Costs'][0], 2023), eur(c['Material_Costs'][1], 2024), eur(c25['COGS'], 2025)],
        ['Gross Margin',
         eur(c['Gross_Margin'][0], 2023), eur(c['Gross_Margin'][1], 2024), eur(c25['Gross_Margin_After_COGS'], 2025)],
        ['  Labor',
         eur(c['Labor_Direct_Indirect'][0], 2023), eur(c['Labor_Direct_Indirect'][1], 2024), eur(c25['Labor_Direct_Indirect'], 2025)],
        ['  Manufacturing',
         eur(c['Manufacturing'][0], 2023), eur(c['Manufacturing'][1], 2024), eur(c25['Manufacturing'], 2025)],
        ['  Delivery / Transport',
         eur(c['Transport_Costs'][0], 2023), eur(c['Transport_Costs'][1], 2024), eur(c25['Delivery'], 2025)],
        ['  Marketing',
         eur(c['Marketing'][0], 2023), eur(c['Marketing'][1], 2024), eur(c25['Marketing'], 2025)],
        ['  G&A',
         eur(c['GA'][0], 2023), eur(c['GA'][1], 2024), eur(c25['GA'], 2025)],
        ['EBITDA',
         eur(c['EBITDA_Normalised'][0], 2023), eur(c['EBITDA_Normalised'][1], 2024), eur(c25['EBITDA'], 2025)],
        ['Net Profit',
         eur(c['Net_Profit'][0], 2023), eur(c['Net_Profit'][1], 2024), eur(c25['Net_Profit'], 2025)],
    ]
    add_table(slide2, margin_left, Cm(3.1), content_width,
              merged_data, [5, 3, 3, 3], font_size=7,
              bold_rows={1, 3, 9, 10})

    # --- Note ---
    add_textbox(slide2, margin_left, Cm(8.5), content_width, Cm(0.6),
        f"2023–2024: Lifebox + Catering only. 2025B includes FruFru relaunch. "
        f"Revenue grew {fmt_pct(growth_rev)} YoY (2023→2024). "
        f"EBITDA negative due to relaunch investment and scaling costs.",
        font_size=6.5, color=LIGHT_TEXT)

    # ==================== SLIDE 3: BUDGET & FORWARD-LOOKING ====================

    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_banner(slide3, margin_left, Cm(0.3), content_width,
                       "FINANCIAL BUDGET — FORWARD LOOKING")

    # --- Budget 2025 by BU ---
    add_textbox(slide3, margin_left, Cm(1.3), int(content_width * 0.5), Cm(0.4),
                "Budget 2025 — By Business Unit",
                font_size=9, bold=True, color=PRIMARY_BLUE)

    bf = fin['bu_frufru_2025']
    bl = fin['bu_lifebox_2025']
    bc = fin['bu_catering_2025']
    total_sales_25 = bf['Sales'] + bl['Sales'] + bc['Sales']

    bu_data = [
        ['Business Unit', 'Sales', 'COGS', 'Gross Margin', '% Total'],
        ['FruFru (Retail)', eur(bf['Sales'], 2025), eur(bf['COGS'], 2025), eur(bf['Gross_Margin'], 2025),
         fmt_pct_abs(bf['Sales'] / total_sales_25)],
        ['Lifebox (Subs)', eur(bl['Sales'], 2025), eur(bl['COGS'], 2025), eur(bl['Gross_Margin'], 2025),
         fmt_pct_abs(bl['Sales'] / total_sales_25)],
        ['Catering (B2B)', eur(bc['Sales'], 2025), eur(bc['COGS'], 2025), eur(bc['Gross_Margin'], 2025),
         fmt_pct_abs(bc['Sales'] / total_sales_25)],
        ['TOTAL', eur(total_sales_25, 2025),
         eur(bf['COGS'] + bl['COGS'] + bc['COGS'], 2025),
         eur(bf['Gross_Margin'] + bl['Gross_Margin'] + bc['Gross_Margin'], 2025),
         '100.0%'],
    ]
    add_table(slide3, margin_left, Cm(1.8), int(content_width * 0.5),
              bu_data, [3, 2, 2, 2, 1.5], font_size=6.5,
              bold_rows={5}, row_height_cm=0.37)

    # --- Budget 2026 Base Case (ALL lines) ---
    # Positioned below the cons25 table: 1.8 + 14×0.37 = ~7.0cm + gap
    b26_title_y = Cm(7.3)
    b26_table_y = Cm(7.7)

    add_textbox(slide3, margin_left, b26_title_y, content_width, Cm(0.4),
                "Budget 2026 — Base Case",
                font_size=9, bold=True, color=PRIMARY_BLUE)

    b26 = fin['budget_2026']
    # EUR conversion helper for Budget 2026 — all values at 2026 exchange rate
    e26 = lambda key: eur(b26.get(key, 0), 2026)
    def fy26(key):
        v = b26.get(key + '_FY', 0)
        return eur(v, 2026) if v > 0 else '—'

    b26_data = [
        ['EUR', 'H1 2026', 'FY 2026E'],
        ['  Sales — Lifebox', e26('Sales_Lifebox_H1'), fy26('Sales_Lifebox')],
        ['  Sales — FruFru', e26('Sales_FruFru_H1'), fy26('Sales_FruFru')],
        ['  Sales — Catering', e26('Sales_Catering_H1'), fy26('Sales_Catering')],
        ['Total Sales', e26('Total_Sales_H1'), eur(b26['Total_Sales'], 2026)],
        ['  Food Cost', e26('Food_Cost_H1'), fy26('Food_Cost')],
        ['  Packaging', e26('Cost_Ambalaj_H1'), fy26('Cost_Ambalaj')],
        ['Total Material Cost', e26('Total_Material_Cost_H1'), eur(b26['Total_Material_Cost_FY'], 2026)],
        ['Material Margin', e26('Material_Margin_H1'), eur(b26['Material_Margin_FY'], 2026)],
        ['  Direct Labor', e26('Salarii_Directe_H1'), fy26('Salarii_Directe')],
        ['  Indirect Labor', e26('Salarii_Indirecte_H1'), fy26('Salarii_Indirecte')],
        ['Total Labor', e26('Total_Salarii_H1'), eur(b26['Total_Salarii_FY'], 2026)],
        ['Total Production Cost', e26('Total_Cost_Productie_H1'), eur(b26['Total_Cost_Productie_FY'], 2026)],
        ['Gross Margin', e26('Gross_Margin_H1'), eur(b26['Gross_Margin_FY'], 2026)],
        ['  Delivery', e26('Delivery_H1'), eur(b26['Delivery_FY'], 2026)],
        ['  Marketing', e26('Marketing_H1'), eur(b26['Marketing_FY'], 2026)],
        ['  G&A', e26('GA_H1'), eur(b26['GA_FY'], 2026)],
        ['EBITDA', e26('EBITDA_Normalised_H1'), eur(b26['EBITDA_Normalised'], 2026)],
    ]
    # Compact rows: 18 × 0.37 = 6.66cm, ending at ~14.36cm
    add_table(slide3, margin_left, b26_table_y, int(content_width * 0.55),
              b26_data, [4, 2.5, 2.5], font_size=6.5,
              bold_rows={4, 7, 8, 11, 13, 17}, row_height_cm=0.37)


    # --- Footer ---
    add_textbox(slide3, margin_left, Cm(17.7), content_width, Cm(0.4),
                "CONFIDENTIAL",
                font_size=6, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)

    # ==================== SLIDE 4: INVESTMENT HIGHLIGHTS & RISKS ====================

    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    add_section_banner(slide4, margin_left, Cm(0.3), content_width, "INVESTMENT HIGHLIGHTS")

    highlights = [
        ("Proven Brand with Unilever Pedigree.",
         "FruFru operated under Unilever for 5 years. Buyback preserves brand equity while restoring agility."),
        ("Integrated Production Platform.",
         "Shared kitchen infrastructure in Bucharest serves three revenue streams with operational leverage."),
        ("Multi-Channel Distribution.",
         "Modern trade (Mega Image, Carrefour, Freshful), D2C (lifebox.ro), B2B catering, "
         "and OTOTO stores (Simiuc's separate retail — distribution synergies)."),
        ("Founder-Operator Alignment.",
         "Simiuc — founder since 2006, PE advisor at Comitis Capital. "
         "Lifebox team (Balaceanu, Scarlat) provides operational continuity."),
        (f"Clear Growth Path.",
         f"Revenue CAGR 2024A→2026E of {fmt_pct(fin['cagr_rev_2024_2026'])} targets "
         f"{fmt_eur(to_eur(rev_2026_bgt, 2026))} with {fmt_pct_abs(ebitda_margin_2026)} EBITDA margin."),
        ("Market Tailwinds.",
         "Romanian healthy food structurally underpenetrated vs. Western Europe. "
         "RTE market ~EUR 1.75B growing 3.7% CAGR. Poland, a comparable CEE economy, "
         "has seen rapid ready meals adoption — Romania expected to follow."),
    ]
    add_bullet_textbox(slide4, margin_left, Cm(1.3), content_width, Cm(5.5),
                       highlights, font_size=8.5, color=DARK_TEXT)

    # --- Product Portfolio ---
    add_section_banner(slide4, margin_left, Cm(7.0), content_width, "PRODUCT PORTFOLIO")

    portfolio_data = [
        ['Brand', 'Channel', 'Products', 'Price Range'],
        ['FruFru', 'Retail (Mega Image, Carrefour, Freshful, Kaufland, OTOTO*)', '~30+ SKUs: salads, soups, bowls, mains', 'RON 15–41'],
        ['Lifebox', 'D2C subscription (daily delivery)', '7 menu types: OptimBox, Vegan, Sport, etc.', 'RON 20.5–30.8/meal'],
        ['B2B Catering', 'Corporate (min 20 portions)', 'Daily lunch, holiday, event platters', 'RON 130–500/platter'],
    ]
    add_table(slide4, margin_left, Cm(8.0), content_width,
              portfolio_data, [2.5, 5, 5.5, 3], font_size=7)

    # --- Competitive Landscape ---
    add_section_banner(slide4, margin_left, Cm(10.2), content_width, "COMPETITIVE LANDSCAPE")

    comp_data = [
        ['Company', 'Segment', 'Revenue', 'Notes'],
        ['Lifebox + FruFru', 'Subscription + Retail', f'{eur(rev_2024, 2024)} (2024)', 'Subject of this memo'],
        ['Salad Box', 'Fast-casual healthy', '~€7.7M (2024)', 'Entered insolvency'],
        ['FoodKit', 'Weekly meal prep', 'Undisclosed', 'Raised EUR 1M (2022)'],
        ["Pep&Pepper", 'Fast-casual healthy', '~EUR 2.8M', '15 locations'],
    ]
    add_table(slide4, margin_left, Cm(11.2), content_width,
              comp_data, [3, 3.5, 3, 5.5], font_size=6.5)

    # --- Footer ---
    add_textbox(slide4, margin_left, Cm(17.7), content_width, Cm(0.4),
                "CONFIDENTIAL",
                font_size=6, color=LIGHT_TEXT, alignment=PP_ALIGN.CENTER)

    # V7 Capital logo — bottom-right of last slide
    if os.path.exists(V7_LOGO):
        # 1080x1350 original → scale to ~2.5cm height
        slide4.shapes.add_picture(V7_LOGO, Cm(24.0), Cm(15.5), height=Cm(2.5))

    # ==================== SAVE ====================

    prs.save(OUTPUT_PPTX)
    print(f"PPTX created: {OUTPUT_PPTX}")
    return OUTPUT_PPTX


if __name__ == "__main__":
    create_info_memo_pptx()
